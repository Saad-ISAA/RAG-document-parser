# parsers/powerpoint_parser.py
"""
PowerPoint parser for PPT and PPTX files with full image and text extraction.
Supports both English and Arabic content with OCR capabilities.
"""

import logging
from pathlib import Path
from typing import List, Optional, Dict, Any
from datetime import datetime
import io
import tempfile

# Import PowerPoint libraries with fallbacks
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.shapes.picture import Picture
    from pptx.shapes.group import GroupShape
    HAS_PYTHON_PPTX = True
except ImportError:
    HAS_PYTHON_PPTX = False
    Presentation = None

# Import image processing for OCR
try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False
    Image = None

try:
    import easyocr
    HAS_EASYOCR = True
except ImportError:
    HAS_EASYOCR = False
    easyocr = None

from models.parse_result import ParseResult, TableData, ImageInfo, DocumentMetadata, FileInfo
from utils.config import ParserConfig

logger = logging.getLogger(__name__)


class PowerPointParser:
    """
    Parser for PowerPoint presentations (PPT and PPTX) with comprehensive
    text and image extraction including OCR support.
    """
    
    def __init__(self, config: ParserConfig):
        """Initialize PowerPoint parser with configuration."""
        self.config = config
        self.ppt_config = config.powerpoint if hasattr(config, 'powerpoint') else None
        self.ocr_config = config.ocr
        
        # Initialize OCR reader for images
        self.ocr_reader = None
        if HAS_EASYOCR and self.ocr_config.enabled:
            try:
                # Include both English and Arabic for requirements
                languages = list(set(self.ocr_config.languages + ['en', 'ar']))
                self.ocr_reader = easyocr.Reader(languages, verbose=False)
                logger.info(f"EasyOCR initialized for PowerPoint with languages: {languages}")
            except Exception as e:
                logger.warning(f"Failed to initialize EasyOCR for PowerPoint: {str(e)}")
        
        logger.info("PowerPoint parser initialized")
    
    def parse(self, file_path: Path, file_info: FileInfo) -> ParseResult:
        """
        Parse a PowerPoint file.
        
        Args:
            file_path: Path to the PowerPoint file
            file_info: File information from detector
            
        Returns:
            ParseResult with extracted content, images, and tables
        """
        result = ParseResult(
            file_path=str(file_path),
            file_type=file_info.mime_type
        )
        
        extension = file_path.suffix.lower()
        
        try:
            if extension == '.pptx':
                return self._parse_pptx(file_path, result)
            elif extension == '.ppt':
                return self._parse_ppt(file_path, result)
            else:
                result.error = f"Unsupported PowerPoint format: {extension}"
                return result
                
        except Exception as e:
            logger.error(f"Error parsing PowerPoint {file_path}: {str(e)}")
            result.error = str(e)
            return result
    
    def _parse_pptx(self, file_path: Path, result: ParseResult) -> ParseResult:
        """Parse PPTX file using python-pptx."""
        if not HAS_PYTHON_PPTX:
            result.error = "python-pptx library not available. Install with: pip install python-pptx"
            return result
        
        try:
            presentation = Presentation(file_path)
            
            # Extract metadata
            result.metadata = self._extract_pptx_metadata(presentation, file_path)
            
            # Extract content from all slides
            slide_contents = []
            all_tables = []
            all_images = []
            
            for slide_index, slide in enumerate(presentation.slides):
                slide_content = f"=== Slide {slide_index + 1} ===\n"
                
                # Extract text from shapes
                slide_text = self._extract_slide_text(slide)
                if slide_text:
                    slide_content += slide_text + "\n"
                
                # Extract tables from slide
                slide_tables = self._extract_slide_tables(slide, slide_index + 1)
                all_tables.extend(slide_tables)
                
                # Add table summaries to slide content
                for table in slide_tables:
                    slide_content += f"\n[Table with {len(table.headers)} columns and {len(table.rows)} rows]\n"
                
                # Extract images from slide
                slide_images = self._extract_slide_images(slide, slide_index + 1)
                all_images.extend(slide_images)
                
                # Add image text to slide content if OCR found text
                for img in slide_images:
                    if img.extracted_text:
                        slide_content += f"\n[Image Text]: {img.extracted_text}\n"
                
                slide_contents.append(slide_content)
            
            # Combine all content
            result.content = '\n\n'.join(slide_contents)
            result.tables = all_tables
            result.images = all_images
            
            result.success = True
            result.parser_used = 'python-pptx'
            
        except Exception as e:
            result.error = f"PPTX parsing failed: {str(e)}"
            logger.error(f"PPTX parsing error: {str(e)}")
        
        return result
    
    def _parse_ppt(self, file_path: Path, result: ParseResult) -> ParseResult:
        """Parse PPT file (legacy format)."""
        # PPT files require conversion or specialized library
        # Try conversion to PPTX first, then parse
        try:
            from utils.document_converter import DocumentConverter
            converter = DocumentConverter(self.config)
            
            # Convert PPT to PPTX
            converted_path = converter.convert_to_pptx(file_path)
            if converted_path and converted_path.exists():
                logger.info(f"Successfully converted {file_path.name} to PPTX for parsing")
                
                # Parse the converted file
                converted_result = self._parse_pptx(converted_path, result)
                converted_result.parser_used = 'ppt-conversion-pptx'
                
                # Cleanup converted file
                try:
                    converted_path.unlink()
                except:
                    pass
                
                return converted_result
            else:
                result.error = "PPT file conversion failed. Unable to parse legacy PPT format."
                
        except ImportError:
            result.error = "PPT parsing requires document converter. Legacy PPT format not directly supported."
        except Exception as e:
            result.error = f"PPT parsing failed: {str(e)}"
        
        return result
    
    def _extract_slide_text(self, slide) -> str:
        """Extract all text from a slide, including text in shapes and groups."""
        text_parts = []
        
        for shape in slide.shapes:
            text = self._extract_shape_text(shape)
            if text:
                text_parts.append(text)
        
        return '\n'.join(text_parts)
    
    def _extract_shape_text(self, shape) -> str:
        """Extract text from a shape, handling different shape types."""
        try:
            # Handle grouped shapes
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_text = []
                for grouped_shape in shape.shapes:
                    text = self._extract_shape_text(grouped_shape)
                    if text:
                        group_text.append(text)
                return '\n'.join(group_text)
            
            # Handle shapes with text frames
            if hasattr(shape, 'text_frame') and shape.text_frame:
                paragraphs = []
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        paragraphs.append(paragraph.text.strip())
                return '\n'.join(paragraphs)
            
            # Handle shapes with text attribute
            if hasattr(shape, 'text') and shape.text:
                return shape.text.strip()
                
        except Exception as e:
            logger.debug(f"Failed to extract text from shape: {str(e)}")
        
        return ""
    
    def _extract_slide_tables(self, slide, slide_number: int) -> List[TableData]:
        """Extract tables from a slide."""
        tables = []
        
        try:
            for shape_index, shape in enumerate(slide.shapes):
                if hasattr(shape, 'table') and shape.table:
                    table_data = self._extract_table_data(shape.table, slide_number, shape_index)
                    if table_data:
                        tables.append(table_data)
        
        except Exception as e:
            logger.warning(f"Table extraction failed on slide {slide_number}: {str(e)}")
        
        return tables
    
    def _extract_table_data(self, table, slide_number: int, table_index: int) -> Optional[TableData]:
        """Extract data from a PowerPoint table."""
        try:
            if not table.rows:
                return None
            
            # Extract headers (first row)
            headers = []
            first_row = table.rows[0]
            for cell in first_row.cells:
                cell_text = ""
                if cell.text_frame:
                    for paragraph in cell.text_frame.paragraphs:
                        cell_text += paragraph.text
                headers.append(cell_text.strip())
            
            # Extract data rows
            rows = []
            for row in table.rows[1:]:  # Skip header row
                row_data = []
                for cell in row.cells:
                    cell_text = ""
                    if cell.text_frame:
                        for paragraph in cell.text_frame.paragraphs:
                            cell_text += paragraph.text
                    row_data.append(cell_text.strip())
                rows.append(row_data)
            
            if headers or rows:
                return TableData(
                    headers=headers,
                    rows=rows,
                    page_number=slide_number,
                    table_index=table_index
                )
        
        except Exception as e:
            logger.warning(f"Failed to extract table data: {str(e)}")
        
        return None
    
    def _extract_slide_images(self, slide, slide_number: int) -> List[ImageInfo]:
        """Extract images from a slide with OCR support."""
        images = []
        
        try:
            image_index = 0
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_info = self._process_image_shape(shape, slide_number, image_index)
                    if image_info:
                        images.append(image_info)
                        image_index += 1
                
                # Handle grouped shapes that might contain images
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    group_images = self._extract_group_images(shape, slide_number, image_index)
                    images.extend(group_images)
                    image_index += len(group_images)
        
        except Exception as e:
            logger.warning(f"Image extraction failed on slide {slide_number}: {str(e)}")
        
        return images
    
    def _process_image_shape(self, shape, slide_number: int, image_index: int) -> Optional[ImageInfo]:
        """Process a single image shape and extract OCR text if enabled."""
        try:
            image_info = ImageInfo(
                image_index=image_index,
                page_number=slide_number
            )
            
            # Get image data
            if hasattr(shape, 'image'):
                image_data = shape.image.blob
                image_info.size_bytes = len(image_data)
                
                # Extract basic image properties
                try:
                    pil_image = Image.open(io.BytesIO(image_data))
                    image_info.width = pil_image.width
                    image_info.height = pil_image.height
                    image_info.format = pil_image.format
                    
                    # Run OCR if enabled and image is large enough
                    if (self.ocr_config.enabled and self.ocr_reader and
                        pil_image.width > 50 and pil_image.height > 30):
                        
                        ocr_text = self._extract_text_from_image(pil_image)
                        if ocr_text:
                            image_info.extracted_text = ocr_text
                            logger.debug(f"OCR extracted {len(ocr_text)} chars from slide {slide_number} image {image_index}")
                
                except Exception as e:
                    logger.debug(f"Failed to process image properties: {str(e)}")
            
            return image_info
            
        except Exception as e:
            logger.debug(f"Failed to process image shape: {str(e)}")
            return None
    
    def _extract_group_images(self, group_shape, slide_number: int, start_index: int) -> List[ImageInfo]:
        """Extract images from grouped shapes."""
        images = []
        image_index = start_index
        
        try:
            for shape in group_shape.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_info = self._process_image_shape(shape, slide_number, image_index)
                    if image_info:
                        images.append(image_info)
                        image_index += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    # Recursive handling of nested groups
                    nested_images = self._extract_group_images(shape, slide_number, image_index)
                    images.extend(nested_images)
                    image_index += len(nested_images)
        
        except Exception as e:
            logger.debug(f"Failed to extract images from group: {str(e)}")
        
        return images
    
    def _extract_text_from_image(self, pil_image: Image.Image) -> Optional[str]:
        """Extract text from image using OCR with Arabic/English support."""
        if not self.ocr_reader:
            return None
        
        try:
            # Preprocess image for better OCR results
            if pil_image.mode not in ['RGB', 'L']:
                pil_image = pil_image.convert('RGB')
            
            # Convert to numpy array for EasyOCR
            import numpy as np
            image_array = np.array(pil_image)
            
            # Run OCR
            results = self.ocr_reader.readtext(image_array)
            
            # Filter by confidence and combine text
            text_parts = []
            for bbox, text, confidence in results:
                if confidence >= self.ocr_config.confidence_threshold:
                    text_parts.append(text.strip())
            
            # Join text parts, handling Arabic RTL properly
            if text_parts:
                # For mixed Arabic/English, preserve line structure
                return '\n'.join(text_parts)
            
        except Exception as e:
            logger.debug(f"OCR failed for PowerPoint image: {str(e)}")
        
        return None
    
    def _extract_pptx_metadata(self, presentation, file_path: Path) -> DocumentMetadata:
        """Extract metadata from PPTX presentation."""
        metadata = DocumentMetadata()
        
        try:
            # Core properties
            if hasattr(presentation, 'core_properties'):
                core_props = presentation.core_properties
                metadata.title = core_props.title
                metadata.author = core_props.author
                metadata.subject = core_props.subject
                metadata.creator = core_props.author
                metadata.creation_date = core_props.created
                metadata.modification_date = core_props.modified
            
            # Slide count
            metadata.page_count = len(presentation.slides)
            
            # Calculate approximate word count
            total_words = 0
            for slide in presentation.slides:
                slide_text = self._extract_slide_text(slide)
                total_words += len(slide_text.split())
            metadata.word_count = total_words
            
            # Custom properties
            metadata.custom_properties = {
                'slide_count': len(presentation.slides),
                'presentation_format': 'PPTX'
            }
            
        except Exception as e:
            logger.warning(f"Failed to extract PPTX metadata: {str(e)}")
        
        return metadata
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported PowerPoint formats."""
        formats = ['pptx']
        if HAS_PYTHON_PPTX:
            formats.append('ppt')  # Through conversion
        return formats
    
    def get_library_info(self) -> Dict[str, Any]:
        """Get information about available PowerPoint libraries."""
        return {
            'python_pptx_available': HAS_PYTHON_PPTX,
            'pil_available': HAS_PIL,
            'ocr_available': HAS_EASYOCR,
            'capabilities': {
                'pptx': {
                    'text_extraction': HAS_PYTHON_PPTX,
                    'table_extraction': HAS_PYTHON_PPTX,
                    'image_extraction': HAS_PYTHON_PPTX,
                    'ocr_support': HAS_EASYOCR,
                    'metadata_extraction': HAS_PYTHON_PPTX
                },
                'ppt': {
                    'text_extraction': True,  # Through conversion
                    'table_extraction': True,  # Through conversion
                    'image_extraction': True,  # Through conversion
                    'ocr_support': HAS_EASYOCR,
                    'metadata_extraction': True  # Through conversion
                }
            },
            'supported_languages': self.ocr_config.languages if self.ocr_config.enabled else [],
            'arabic_support': 'ar' in self.ocr_config.languages if self.ocr_config.enabled else False
        }